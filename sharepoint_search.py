"""
title: SharePoint Search
author: Eshan Shah
version: 3.4.0
license: MIT
description: Search SharePoint files, pages, and list items from Open Web UI using Microsoft Graph Search API with delegated (per-user) authentication.
requirements: PyPDF2, python-docx, python-pptx, openpyxl
"""

import base64
import hashlib
import io
import json
import re
import sqlite3
import time
import asyncio
import math
from typing import Optional
from cryptography.fernet import Fernet, InvalidToken
from pydantic import BaseModel, Field

import aiohttp


class Tools:
    class Valves(BaseModel):
        azure_client_id: str = Field(
            default="",
            description="Azure AD app registration client ID"
        )
        azure_tenant_id: str = Field(
            default="common",
            description="Azure AD tenant ID, or 'common' for multi-tenant"
        )
        encryption_key: str = Field(
            default="",
            description="Secret key for encrypting stored tokens. Set any random string (e.g. a UUID). If empty, tokens are stored in plaintext."
        )
        sharepoint_site_ids: str = Field(
            default="",
            description="Semicolon-separated SharePoint site IDs to scope search (e.g. 'host.sharepoint.com,guid1,guid1;host.sharepoint.com,guid2,guid2'). Leave empty to search all sites the user has access to."
        )
        max_results: int = Field(
            default=10,
            description="Maximum number of search results to return"
        )
        max_content_chars: int = Field(
            default=50000,
            description="Maximum characters of extracted text to return per file. Content beyond this limit is truncated."
        )
        max_download_mb: int = Field(
            default=50,
            description="Maximum file size (in MB) the connector will download from SharePoint. Files larger than this are rejected before download to protect memory."
        )
        child_chunk_size: int = Field(
            default=300,
            description="Characters per child chunk scored by BM25. Smaller = more precise matching. Parent sections returned to the LLM are children_per_parent × this size (~1500 chars at defaults)."
        )
        children_per_parent: int = Field(
            default=5,
            description="Number of child chunks grouped into one parent section returned to the LLM. Default 5 × 300 chars = ~1500 chars per section."
        )
        top_k_chunks: int = Field(
            default=8,
            description="Number of top-scoring parent sections to return when query-based retrieval is active. Each section is ~children_per_parent × child_chunk_size chars."
        )
        min_bm25_score: float = Field(
            default=0.01,
            description="Minimum BM25 score for a child chunk to count as a keyword match. If all chunks fall below this, falls back to returning the first 3 parent sections."
        )

    GRAPH_SCOPE = "https://graph.microsoft.com/Sites.Read.All https://graph.microsoft.com/User.Read offline_access"

    def __init__(self):
        self.valves = self.Valves()
        self._db_path = "/app/backend/data/sp_search.db"
        self._fernet: Optional[Fernet] = None
        self._init_db()

    def _get_fernet(self) -> Optional[Fernet]:
        """Derive a Fernet key from the Valve-configured encryption_key."""
        raw = self.valves.encryption_key.strip()
        if not raw:
            return None
        key_bytes = hashlib.sha256(raw.encode()).digest()
        fernet_key = base64.urlsafe_b64encode(key_bytes)
        return Fernet(fernet_key)

    def _encrypt(self, plaintext: str) -> str:
        f = self._get_fernet()
        if not f:
            return plaintext
        return f.encrypt(plaintext.encode()).decode()

    def _decrypt(self, ciphertext: str) -> str:
        f = self._get_fernet()
        if not f:
            return ciphertext
        try:
            return f.decrypt(ciphertext.encode()).decode()
        except InvalidToken:
            # Token was stored with a different key or unencrypted — force re-auth
            return ""

    def _init_db(self):
        conn = sqlite3.connect(self._db_path)
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("""
            CREATE TABLE IF NOT EXISTS user_auth (
                owui_user_id   TEXT PRIMARY KEY,
                access_token   TEXT,
                refresh_token  TEXT,
                token_expiry   INTEGER DEFAULT 0
            )
        """)
        conn.commit()
        conn.close()

    def _ensure_table(self, conn):
        conn.execute("""
            CREATE TABLE IF NOT EXISTS user_auth (
                owui_user_id   TEXT PRIMARY KEY,
                access_token   TEXT,
                refresh_token  TEXT,
                token_expiry   INTEGER DEFAULT 0
            )
        """)

    def _get_user_auth(self, user_id: str) -> Optional[dict]:
        conn = sqlite3.connect(self._db_path)
        self._ensure_table(conn)
        conn.row_factory = sqlite3.Row
        row = conn.execute(
            "SELECT * FROM user_auth WHERE owui_user_id = ?", (user_id,)
        ).fetchone()
        conn.close()
        if row is None:
            return None
        auth = dict(row)
        # Decrypt token fields
        if auth.get("access_token"):
            auth["access_token"] = self._decrypt(auth["access_token"])
        if auth.get("refresh_token"):
            auth["refresh_token"] = self._decrypt(auth["refresh_token"])
        return auth

    def _save_user_auth(self, user_id: str, **kwargs):
        # Encrypt token fields before storing
        store = dict(kwargs)
        if "access_token" in store and store["access_token"]:
            store["access_token"] = self._encrypt(store["access_token"])
        if "refresh_token" in store and store["refresh_token"]:
            store["refresh_token"] = self._encrypt(store["refresh_token"])

        conn = sqlite3.connect(self._db_path)
        self._ensure_table(conn)
        existing = conn.execute(
            "SELECT 1 FROM user_auth WHERE owui_user_id = ?", (user_id,)
        ).fetchone()
        if existing is None:
            fields = {"owui_user_id": user_id, **store}
            cols = ", ".join(fields.keys())
            placeholders = ", ".join(["?"] * len(fields))
            conn.execute(f"INSERT INTO user_auth ({cols}) VALUES ({placeholders})", list(fields.values()))
        else:
            set_clause = ", ".join(f"{k} = ?" for k in store.keys())
            conn.execute(
                f"UPDATE user_auth SET {set_clause} WHERE owui_user_id = ?",
                list(store.values()) + [user_id]
            )
        conn.commit()
        conn.close()

    async def _emit(self, emitter, description: str, done: bool = False):
        if emitter:
            await emitter({
                "type": "status",
                "data": {"description": description, "done": done}
            })

    async def _get_valid_token(self, user_id: str) -> Optional[str]:
        """Return a valid access token for the user, refreshing if needed."""
        auth = self._get_user_auth(user_id)
        if auth is None:
            return None

        # Token still fresh
        if auth["token_expiry"] > int(time.time()) + 60:
            return auth["access_token"]

        # Refresh
        if not auth.get("refresh_token"):
            return None

        tenant = self.valves.azure_tenant_id
        client_id = self.valves.azure_client_id

        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=30)) as session:
            async with session.post(
                f"https://login.microsoftonline.com/{tenant}/oauth2/v2.0/token",
                data={
                    "grant_type": "refresh_token",
                    "client_id": client_id,
                    "refresh_token": auth["refresh_token"],
                    "scope": self.GRAPH_SCOPE,
                },
            ) as resp:
                if resp.status != 200:
                    return None
                token_data = await resp.json()

        if "access_token" not in token_data:
            return None

        self._save_user_auth(
            user_id,
            access_token=token_data["access_token"],
            refresh_token=token_data.get("refresh_token", auth["refresh_token"]),
            token_expiry=int(time.time()) + token_data.get("expires_in", 3600),
        )
        return token_data["access_token"]

    def _get_site_filter(self) -> str:
        """Build a KQL site filter from configured site IDs to scope search to specific sites."""
        raw = self.valves.sharepoint_site_ids.strip()
        if not raw:
            return ""
        site_ids = [s.strip() for s in raw.split(";") if s.strip()]
        # Use KQL siteId filter: (siteId:guid1 OR siteId:guid2)
        # Extract the middle GUID from "host,siteGuid,webGuid" format
        guids = []
        for sid in site_ids:
            parts = sid.split(",")
            if len(parts) >= 2:
                guids.append(parts[1].strip())
        if not guids:
            return ""
        if len(guids) == 1:
            return f"siteId:{guids[0]}"
        return "(" + " OR ".join(f"siteId:{g}" for g in guids) + ")"

    async def connect_sharepoint(
        self,
        __user__: dict = {},
        __event_emitter__=None,
    ) -> str:
        """
        Connect your Microsoft account to enable SharePoint search. This starts a device code sign-in flow.

        :return: Instructions for signing in, or confirmation of connection.
        """
        user_id = __user__.get("id", "")
        if not user_id:
            return "Error: Could not identify user."

        if not self.valves.azure_client_id:
            return "Error: SharePoint search is not configured. An admin must set the Azure Client ID in Valves."

        # Check if already connected
        auth = self._get_user_auth(user_id)
        if auth and auth.get("refresh_token"):
            token = await self._get_valid_token(user_id)
            if token:
                return "You're already connected to SharePoint. Use 'disconnect SharePoint' if you want to reconnect."

        tenant = self.valves.azure_tenant_id
        client_id = self.valves.azure_client_id

        await self._emit(__event_emitter__, "Requesting device code from Microsoft...")

        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=30)) as session:
            async with session.post(
                f"https://login.microsoftonline.com/{tenant}/oauth2/v2.0/devicecode",
                data={"client_id": client_id, "scope": self.GRAPH_SCOPE},
            ) as resp:
                if resp.status != 200:
                    error_text = await resp.text()
                    await self._emit(__event_emitter__, "Failed to get device code.", done=True)
                    return f"Error requesting device code: {error_text}"
                device_data = await resp.json()

        user_code = device_data["user_code"]
        verification_uri = device_data["verification_uri"]
        device_code = device_data["device_code"]
        expires_in = device_data.get("expires_in", 900)
        interval = device_data.get("interval", 5)

        await self._emit(__event_emitter__, f"Waiting for sign-in...")

        # Send clickable link and code to the chat
        if __event_emitter__:
            await __event_emitter__({
                "type": "message",
                "data": {"content": f"\n\n**Sign in to connect SharePoint:**\n\n1. Click here: [{verification_uri}]({verification_uri})\n2. Enter code: **{user_code}**\n3. Sign in with your Microsoft work account\n\n*Waiting for you to complete sign-in...*\n\n"}
            })

        # Poll for token
        token_url = f"https://login.microsoftonline.com/{tenant}/oauth2/v2.0/token"
        deadline = time.time() + expires_in

        while time.time() < deadline:
            await asyncio.sleep(interval)
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=30)) as session:
                async with session.post(
                    token_url,
                    data={
                        "grant_type": "urn:ietf:params:oauth:grant-type:device_code",
                        "client_id": client_id,
                        "device_code": device_code,
                    },
                ) as resp:
                    token_data = await resp.json()

            if "access_token" in token_data:
                self._save_user_auth(
                    user_id,
                    access_token=token_data["access_token"],
                    refresh_token=token_data.get("refresh_token", ""),
                    token_expiry=int(time.time()) + token_data.get("expires_in", 3600),
                )
                await self._emit(__event_emitter__, "Connected to SharePoint!", done=True)
                return "Successfully connected to Microsoft. You can now search SharePoint."

            error = token_data.get("error", "")
            if error == "authorization_pending":
                continue
            elif error == "slow_down":
                interval += 5
                continue
            else:
                await self._emit(__event_emitter__, "Authentication failed.", done=True)
                return f"Authentication failed: {token_data.get('error_description', error)}"

        await self._emit(__event_emitter__, "Authentication timed out.", done=True)
        return "Device code expired. Please try connecting again."

    async def disconnect_sharepoint(
        self,
        __user__: dict = {},
        __event_emitter__=None,
    ) -> str:
        """
        Disconnect your Microsoft account from SharePoint search.

        :return: Confirmation message.
        """
        user_id = __user__.get("id", "")
        if not user_id:
            return "Error: Could not identify user."

        conn = sqlite3.connect(self._db_path)
        self._ensure_table(conn)
        conn.execute("DELETE FROM user_auth WHERE owui_user_id = ?", (user_id,))
        conn.commit()
        conn.close()
        return "SharePoint disconnected. Your tokens have been cleared."

    async def search_sharepoint(
        self,
        query: str,
        __user__: dict = {},
        __event_emitter__=None,
    ) -> str:
        """
        Search SharePoint for files, pages, and list items. Always use this tool first when the user asks about SharePoint content. Each result includes a Resource ID. When the user's question needs information from inside a file (what it says, a summary, specific details), do not stop at the search results — call get_sharepoint_file_content on the single most relevant result to read its full content, then answer. Only read additional files if the first file does not contain enough information and the user explicitly needs more. Only present the bare list without reading when the user just wants to find or locate files.

        :param query: The search query string.
        :return: Formatted search results with titles, snippets, URLs, and resource IDs.
        """
        user_id = __user__.get("id", "")
        if not user_id:
            return "Error: Could not identify user."

        if not self.valves.azure_client_id:
            return "Error: SharePoint search is not configured. An admin must set the Azure Client ID in Valves."

        # Check auth
        auth = self._get_user_auth(user_id)
        if auth is None or not auth.get("refresh_token"):
            return "You haven't connected your Microsoft account yet. Say 'connect SharePoint' to get started."

        await self._emit(__event_emitter__, "Searching SharePoint...")

        token = await self._get_valid_token(user_id)
        if not token:
            return "Your Microsoft session has expired. Please reconnect with 'connect SharePoint'."

        # Build search request — use KQL site filter to scope to configured sites
        site_filter = self._get_site_filter()
        search_query = f"{query} {site_filter}".strip() if site_filter else query

        search_body = {
            "requests": [
                {
                    "entityTypes": ["driveItem", "listItem"],
                    "query": {"queryString": search_query},
                    "from": 0,
                    "size": self.valves.max_results,
                }
            ]
        }

        # Execute search
        headers = {
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        }

        try:
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=30)) as session:
                async with session.post(
                    "https://graph.microsoft.com/v1.0/search/query",
                    headers=headers,
                    json=search_body,
                ) as resp:
                    if resp.status != 200:
                        error_text = await resp.text()
                        await self._emit(__event_emitter__, f"Search failed ({resp.status})", done=True)
                        return f"Error: Search API returned {resp.status}: {error_text}"
                    data = await resp.json()
        except Exception as e:
            await self._emit(__event_emitter__, f"Request failed: {e}", done=True)
            return f"Error: Failed to reach Microsoft Graph API: {e}"

        # Parse results
        results = []
        for connection in data.get("value", []):
            for hit_container in connection.get("hitsContainers", []):
                if hit_container.get("total", 0) == 0:
                    continue
                for hit in hit_container.get("hits", []):
                    resource = hit.get("resource", {})
                    title = resource.get("name") or resource.get("displayName") or "Untitled"
                    summary = hit.get("summary", "")
                    url = resource.get("webUrl", "")

                    # Extract driveId and itemId for content retrieval
                    hit_id = hit.get("hitId", "")
                    drive_id = resource.get("parentReference", {}).get("driveId", "")
                    item_id = resource.get("id", "")
                    # Fallback: parse from hitId (format: sites/{siteId}/drives/{driveId}/items/{itemId})
                    if not drive_id and hit_id:
                        m = re.search(r"drives/([^/]+)/items/([^/]+)", hit_id)
                        if m:
                            drive_id, item_id = m.group(1), m.group(2)

                    resource_id = f"{drive_id}/{item_id}" if drive_id and item_id else ""

                    results.append({
                        "title": title,
                        "snippet": summary,
                        "url": url,
                        "resource_id": resource_id,
                    })

        if not results:
            await self._emit(__event_emitter__, "No results found", done=True)
            return f"No results found for: {query}"

        # Format output — embed links in titles so the LLM preserves them
        output_lines = [f"Found {len(results)} result(s) for: {query}\n"]
        output_lines.append("IMPORTANT: Always include the links below in your response to the user.\n")
        for i, r in enumerate(results, 1):
            if r["url"]:
                safe_url = r["url"].replace(" ", "%20")
                _office_exts = {'.docx', '.doc', '.docm', '.xlsx', '.xls', '.xlsm', '.pptx', '.ppt', '.pptm', '.one'}
                _base = safe_url.split('?')[0].lower()
                _ext = '.' + _base.rsplit('.', 1)[-1] if '.' in _base.rsplit('/', 1)[-1] else ''
                if _ext in _office_exts:
                    safe_url += '?web=1' if '?' not in safe_url else '&web=1'
                output_lines.append(f"{i}. [{r['title']}]({safe_url})")
            else:
                output_lines.append(f"{i}. **{r['title']}**")
            if r["snippet"]:
                output_lines.append(f"   {r['snippet']}")
            if r["resource_id"]:
                output_lines.append(f"   Resource ID: `{r['resource_id']}`")
            output_lines.append("")

        # Prompt the LLM to offer content retrieval
        has_resource_ids = any(r["resource_id"] for r in results)
        if has_resource_ids:
            output_lines.append("---")
            output_lines.append("If answering the user's question needs the contents of a file, call get_sharepoint_file_content on the single most relevant file using its Resource ID AND the user's original question as the `query` argument. Only read one file unless the user explicitly asks for more. If the Resource ID is unavailable, use the file name in the `name` argument instead.")

        await self._emit(__event_emitter__, f"Found {len(results)} result(s)", done=True)
        return "\n".join(output_lines)

    # Text-extractable file extensions
    _TEXT_EXTS = {'.txt', '.csv', '.md', '.json', '.xml', '.html', '.htm', '.log', '.yaml', '.yml', '.ini', '.cfg', '.conf', '.py', '.js', '.ts', '.css', '.sql', '.sh', '.bat', '.ps1'}

    @staticmethod
    def _extract_text_from_docx(file_bytes: bytes, max_chars: int = 0) -> str:
        """
        Extract text from a .docx file. Bails out once `max_chars` is reached
        (0 = no cap). Saves a ton of work + memory on huge docs.
        """
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        parts = []
        running = 0
        budget = max_chars * 2 if max_chars else None  # leave headroom for final truncation

        def _add(s: str) -> bool:
            """Append `s`. Return False if we've blown the budget (caller should stop)."""
            nonlocal running
            parts.append(s)
            running += len(s) + 2  # +2 for the join separator
            return budget is None or running < budget

        # Headers and footers
        for section in doc.sections:
            for header in (section.header, section.first_page_header):
                if header and header.is_linked_to_previous is False:
                    for p in header.paragraphs:
                        if p.text.strip() and not _add(p.text.strip()):
                            return "\n\n".join(parts) + "\n\n[... extraction stopped early — content cap reached]"
            for footer in (section.footer, section.first_page_footer):
                if footer and footer.is_linked_to_previous is False:
                    for p in footer.paragraphs:
                        if p.text.strip() and not _add(p.text.strip()):
                            return "\n\n".join(parts) + "\n\n[... extraction stopped early — content cap reached]"

        # Body paragraphs
        for p in doc.paragraphs:
            if p.text.strip() and not _add(p.text.strip()):
                return "\n\n".join(parts) + "\n\n[... extraction stopped early — content cap reached]"

        # Tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells and not _add("\t".join(cells)):
                    return "\n\n".join(parts) + "\n\n[... extraction stopped early — content cap reached]"

        return "\n\n".join(parts)

    @staticmethod
    def _extract_text_from_pptx(file_bytes: bytes, max_chars: int = 0) -> str:
        """Extract text from a .pptx file. Stops once `max_chars` is reached."""
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        running = 0
        budget = max_chars * 2 if max_chars else None

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                header = f"--- Slide {slide_num} ---"
                lines.append(header)
                lines.extend(slide_texts)
                running += len(header) + sum(len(t) + 1 for t in slide_texts)
                if budget is not None and running >= budget:
                    lines.append("[... extraction stopped early — content cap reached]")
                    break
        return "\n".join(lines)

    @staticmethod
    def _extract_text_from_xlsx(file_bytes: bytes, max_chars: int = 0) -> str:
        """Extract text from an .xlsx file. Stops once `max_chars` is reached."""
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        lines = []
        running = 0
        budget = max_chars * 2 if max_chars else None
        stopped = False

        try:
            for sheet_name in wb.sheetnames:
                if stopped:
                    break
                ws = wb[sheet_name]
                header = f"--- Sheet: {sheet_name} ---"
                lines.append(header)
                running += len(header) + 1
                for row in ws.iter_rows(values_only=True):
                    cell_values = [str(c) if c is not None else "" for c in row]
                    if any(v for v in cell_values):
                        joined = "\t".join(cell_values)
                        lines.append(joined)
                        running += len(joined) + 1
                        if budget is not None and running >= budget:
                            lines.append("[... extraction stopped early — content cap reached]")
                            stopped = True
                            break
        finally:
            wb.close()
        return "\n".join(lines)

    @staticmethod
    def _extract_text_from_pdf(file_bytes: bytes, max_chars: int = 0) -> str:
        """Extract text from a PDF file. Stops once `max_chars` is reached (saves a lot on huge PDFs)."""
        budget = max_chars * 2 if max_chars else None

        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            pages = []
            running = 0
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    chunk = f"--- Page {i} ---\n{text.strip()}"
                    pages.append(chunk)
                    running += len(chunk) + 2
                    if budget is not None and running >= budget:
                        pages.append("[... extraction stopped early — content cap reached]")
                        break
            return "\n\n".join(pages)
        except ImportError:
            pass
        try:
            import pdfplumber
            pages = []
            running = 0
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        chunk = f"--- Page {i} ---\n{text.strip()}"
                        pages.append(chunk)
                        running += len(chunk) + 2
                        if budget is not None and running >= budget:
                            pages.append("[... extraction stopped early — content cap reached]")
                            break
            return "\n\n".join(pages)
        except ImportError:
            pass
        raise ImportError("No PDF library available (need PyPDF2 or pdfplumber)")

    @staticmethod
    def _chunk_text_parent_child(text: str, child_size: int = 300, children_per_parent: int = 5) -> tuple:
        """
        Build a two-level chunk hierarchy for parent-child retrieval.
        Child chunks are scored by BM25; matched children are expanded to their parent
        section before being returned to the LLM, preserving context without flooding tokens.

        Returns (child_chunks, parent_chunks, child_to_parent):
          child_chunks: [(idx, text), ...] — small chunks for BM25 scoring
          parent_chunks: [(idx, text), ...] — larger sections returned to the LLM
          child_to_parent: [parent_idx, ...] — maps child index → parent index
        """
        _SENT_SEPS = ('. ', '! ', '? ', '.\n', '!\n', '?\n')

        # Step 1: sentence-aligned units no larger than child_size
        units = []
        for para in re.split(r'\n\n+', text):
            para = para.strip()
            if not para:
                continue
            if len(para) <= child_size:
                units.append(para)
                continue
            rem = para
            while len(rem) > child_size:
                cut = child_size
                for sep in _SENT_SEPS:
                    idx = rem.rfind(sep, 0, child_size)
                    if idx != -1:
                        cut = idx + len(sep)
                        break
                units.append(rem[:cut].strip())
                rem = rem[cut:].strip()
            if rem:
                units.append(rem)

        # Step 2: pack units into child chunks (no overlap — clean parent grouping)
        child_texts = []
        buf, buf_len = [], 0
        for unit in units:
            if buf and buf_len + len(unit) + 2 > child_size:
                child_texts.append('\n\n'.join(buf).strip())
                buf, buf_len = [], 0
            buf.append(unit)
            buf_len += len(unit) + 2
        if buf:
            chunk = '\n\n'.join(buf).strip()
            if chunk:
                child_texts.append(chunk)

        # Step 3: group consecutive children into parents
        child_to_parent = [i // children_per_parent for i in range(len(child_texts))]
        n_parents = (len(child_texts) + children_per_parent - 1) // children_per_parent if child_texts else 0
        parent_texts = []
        for p in range(n_parents):
            start = p * children_per_parent
            end = min(start + children_per_parent, len(child_texts))
            parent_texts.append('\n\n'.join(child_texts[start:end]))

        return (
            list(enumerate(child_texts)),
            list(enumerate(parent_texts)),
            child_to_parent,
        )

    @staticmethod
    def _bm25_scores(chunks: list, query: str) -> list:
        """Return BM25 relevance score for each chunk string against query. k1=1.5, b=0.75."""
        _STOPWORDS = {
            'the','a','an','is','are','was','were','be','been','being',
            'have','has','had','do','does','did','will','would','could',
            'should','may','might','must','shall','can','of','in','to',
            'and','or','but','if','for','on','at','by','from','with',
            'as','into','it','its','this','that','these','those','i',
            'you','he','she','we','they','me','him','her','us','them',
            'my','your','his','our','their','what','which','who','when',
            'where','why','how','not','no','s','t',
        }

        def tokenize(text):
            return [w for w in re.split(r'[^a-z0-9]+', text.lower())
                    if w and w not in _STOPWORDS]

        if not chunks:
            return []

        tokenized = [tokenize(c) for c in chunks]
        query_terms = tokenize(query)

        if not query_terms:
            return [0.0] * len(chunks)

        N = len(chunks)
        avgdl = sum(len(tc) for tc in tokenized) / N
        k1, b = 1.5, 0.75

        idf = {}
        for term in set(query_terms):
            df = sum(1 for tc in tokenized if term in tc)
            idf[term] = math.log((N - df + 0.5) / (df + 0.5) + 1)

        scores = []
        for tc in tokenized:
            dl = len(tc)
            tf_map: dict = {}
            for w in tc:
                tf_map[w] = tf_map.get(w, 0) + 1
            score = 0.0
            denom_base = k1 * (1 - b + b * dl / (avgdl or 1))
            for term in query_terms:
                tf = tf_map.get(term, 0)
                score += idf.get(term, 0) * (tf * (k1 + 1)) / (tf + denom_base)
            scores.append(score)

        return scores

    async def _download_and_extract(self, drive_id: str, item_id: str, token: str, emitter=None) -> str:
        """Internal helper: download a file from Graph API and extract its text content."""
        headers = {"Authorization": f"Bearer {token}"}

        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=60)) as session:
            # Get file metadata
            async with session.get(
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}",
                headers=headers,
            ) as resp:
                if resp.status != 200:
                    error_text = await resp.text()
                    return f"[Error: Could not get file metadata ({resp.status}): {error_text}]"
                metadata = await resp.json()

            file_name = metadata.get("name", "unknown")
            file_size = metadata.get("size", 0)
            max_size = self.valves.max_download_mb * 1024 * 1024

            if file_size > max_size:
                return f"[Error: File '{file_name}' is {file_size // (1024*1024)} MB, exceeds {self.valves.max_download_mb} MB limit]"

            await self._emit(emitter, f"Downloading {file_name}...")

            # Download file content
            async with session.get(
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}/content",
                headers=headers,
            ) as resp:
                if resp.status not in (200, 302):
                    error_text = await resp.text()
                    return f"[Error: Could not download '{file_name}' ({resp.status}): {error_text}]"

                # Stream the download in chunks with a defensive size guard
                # (in case Graph metadata lied about size or the file grew between calls)
                chunks = []
                downloaded = 0
                async for chunk in resp.content.iter_chunked(1024 * 1024):  # 1 MB chunks
                    downloaded += len(chunk)
                    if downloaded > max_size:
                        return (
                            f"[Error: File '{file_name}' exceeded {self.valves.max_download_mb} MB "
                            f"during download — aborted to protect memory]"
                        )
                    chunks.append(chunk)
                    # Emit progress every ~5 MB
                    if len(chunks) % 5 == 0:
                        await self._emit(emitter, f"Downloading {file_name}... ({downloaded // (1024*1024)} MB)")
                file_bytes = b"".join(chunks)
                del chunks  # free the list before extraction

        # Determine file type and extract text — extraction is sync + CPU-heavy,
        # so we run it in a worker thread to keep the OWUI event loop responsive.
        ext = ('.' + file_name.rsplit('.', 1)[-1].lower()) if '.' in file_name else ''
        max_chars = self.valves.max_content_chars

        await self._emit(emitter, f"Extracting text from {file_name}...")

        try:
            if ext in self._TEXT_EXTS:
                text = file_bytes.decode('utf-8', errors='replace')
            elif ext in ('.docx', '.docm'):
                text = await asyncio.to_thread(self._extract_text_from_docx, file_bytes, max_chars)
            elif ext in ('.pptx', '.pptm'):
                text = await asyncio.to_thread(self._extract_text_from_pptx, file_bytes, max_chars)
            elif ext in ('.xlsx', '.xlsm'):
                text = await asyncio.to_thread(self._extract_text_from_xlsx, file_bytes, max_chars)
            elif ext == '.pdf':
                text = await asyncio.to_thread(self._extract_text_from_pdf, file_bytes, max_chars)
            else:
                return f"[Error: Cannot extract text from '{file_name}' (unsupported format: {ext or 'unknown'})]"
        except ImportError as e:
            return f"[Error: Cannot extract text from '{file_name}': {e}]"
        except MemoryError:
            return (
                f"[Error: Ran out of memory extracting '{file_name}'. "
                f"Lower max_download_mb (currently {self.valves.max_download_mb} MB) "
                f"or give the OWUI container more memory.]"
            )
        except Exception as e:
            return f"[Error: Failed to extract text from '{file_name}': {e}]"
        finally:
            del file_bytes  # free the raw download as soon as extraction finishes

        if not text or not text.strip():
            return f"[File '{file_name}' appears empty or contains no extractable text]"

        text = text.strip()
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n[... truncated — showing {max_chars:,} of {len(text):,} characters]"

        return f"**Content of: {file_name}**\n\n{text}"

    async def _search_for_resource_id(self, query: str, token: str):
        """Run a Graph search and return (resource_id, title) of the best match, or ("", "")."""
        site_filter = self._get_site_filter()
        search_query = f"{query} {site_filter}".strip() if site_filter else query
        search_body = {
            "requests": [
                {
                    "entityTypes": ["driveItem"],
                    "query": {"queryString": search_query},
                    "from": 0,
                    "size": 5,
                }
            ]
        }
        headers = {
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        }
        try:
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=30)) as session:
                async with session.post(
                    "https://graph.microsoft.com/v1.0/search/query",
                    headers=headers,
                    json=search_body,
                ) as resp:
                    if resp.status != 200:
                        return "", ""
                    data = await resp.json()
        except Exception:
            return "", ""

        for connection in data.get("value", []):
            for hit_container in connection.get("hitsContainers", []):
                if hit_container.get("total", 0) == 0:
                    continue
                for hit in hit_container.get("hits", []):
                    resource = hit.get("resource", {})
                    title = resource.get("name") or resource.get("displayName") or "Untitled"
                    hit_id = hit.get("hitId", "")
                    drive_id = resource.get("parentReference", {}).get("driveId", "")
                    item_id = resource.get("id", "")
                    if not drive_id and hit_id:
                        m = re.search(r"drives/([^/]+)/items/([^/]+)", hit_id)
                        if m:
                            drive_id, item_id = m.group(1), m.group(2)
                    if drive_id and item_id:
                        return f"{drive_id}/{item_id}", title
        return "", ""

    async def get_sharepoint_file_content(
        self,
        resource_id: str = "",
        name: str = "",
        query: str = "",
        __user__: dict = {},
        __event_emitter__=None,
    ) -> str:
        """
        Read the text content of a SharePoint file, returning only the sections most relevant to the user's question. Call this whenever you need to answer a question from inside a document — summaries, specific details, what a file says.

        ALWAYS pass the user's original question in the `query` argument. This triggers focused parent-child retrieval: the file is chunked, BM25-scored against the query, and only the highest-scoring sections (~1500 chars each) are returned. Without `query`, the raw full document is returned and will overflow the context window for large files.

        Prefer `resource_id` from search_sharepoint results. If unavailable, pass the file name or keywords in `name` and this tool will locate the file.

        :param resource_id: Optional. The resource ID from search results (format: driveId/itemId).
        :param name: Optional. A file name or short search query to locate the file when no resource_id is available.
        :param query: REQUIRED for focused retrieval. The user's original question or topic. Omit only when you explicitly need the raw full text.
        :return: The most relevant sections of the file focused on the query, or the full text if no query is provided.
        """
        user_id = __user__.get("id", "")
        if not user_id:
            return "Error: Could not identify user."

        if not self.valves.azure_client_id:
            return "Error: SharePoint search is not configured."

        # Auth check (needed before we can resolve a name via search)
        auth = self._get_user_auth(user_id)
        if auth is None or not auth.get("refresh_token"):
            return "You haven't connected your Microsoft account yet. Say 'connect SharePoint' to get started."

        token = await self._get_valid_token(user_id)
        if not token:
            return "Your Microsoft session has expired. Please reconnect with 'connect SharePoint'."

        # Resolve which file to read. Prefer an explicit resource_id; otherwise fall
        # back to locating the file by name/query via a Graph search. This keeps the
        # tool usable even when the resource_id from a previous search is no longer in
        # the model's context (e.g. after a new turn or after the orchestrator re-routes).
        resource_id = (resource_id or "").strip()
        parts = resource_id.split("/")
        if not (len(parts) == 2 and all(parts)):
            if name and name.strip():
                await self._emit(__event_emitter__, f"Locating '{name.strip()}'...")
                resolved, _ = await self._search_for_resource_id(name.strip(), token)
                if not resolved:
                    await self._emit(__event_emitter__, "File not found", done=True)
                    return (f"Error: Could not find a SharePoint file matching '{name.strip()}'. "
                            "Try search_sharepoint first, then call this tool with the Resource ID.")
                resource_id = resolved
                parts = resource_id.split("/")
            else:
                return ("Error: No file specified. Provide a resource_id (format 'driveId/itemId') "
                        "from search_sharepoint results, or pass a file name in the `name` argument.")
        drive_id, item_id = parts

        await self._emit(__event_emitter__, "Downloading file from SharePoint...")

        try:
            result = await self._download_and_extract(drive_id, item_id, token, __event_emitter__)
        except Exception as e:
            await self._emit(__event_emitter__, f"Request failed: {e}", done=True)
            return f"Error: Failed to download file: {e}"

        # If a query was provided, apply parent-child BM25 retrieval instead of returning full text.
        # BM25 scores small child chunks for precision; matched children expand to their larger
        # parent section before being returned, preserving context without flooding the LLM.
        # Error strings from _download_and_extract (e.g. "[Error: ...]") are returned as-is.
        if query and query.strip() and result.startswith("**Content of:"):
            header_end = result.find("\n\n")
            if header_end != -1:
                file_header = result[:header_end]
                raw_text = result[header_end + 2:]
            else:
                file_header = ""
                raw_text = result

            child_chunks, parent_chunks, child_to_parent = self._chunk_text_parent_child(
                raw_text,
                child_size=self.valves.child_chunk_size,
                children_per_parent=self.valves.children_per_parent,
            )
            n_children = len(child_chunks)
            n_parents = len(parent_chunks)

            if n_children == 0:
                await self._emit(__event_emitter__, "Done", done=True)
                return result

            # BM25-score child chunks against query
            child_texts = [c for _, c in child_chunks]
            scores = self._bm25_scores(child_texts, query.strip())

            k = min(self.valves.top_k_chunks, n_children)
            ranked = sorted(range(n_children), key=lambda i: scores[i], reverse=True)[:k]
            best_score = scores[ranked[0]] if ranked else 0.0
            any_match = best_score >= self.valves.min_bm25_score

            if any_match:
                # Map top-k children → their parent sections, dedup, preserve doc order
                parent_indices = sorted(set(child_to_parent[i] for i in ranked))
            else:
                # No keyword matches — fall back to opening sections (intro/summary usually first)
                parent_indices = list(range(min(3, n_parents)))

            fname = file_header.replace("**Content of: ", "").replace("**", "").strip()
            n_returned = len(parent_indices)
            lines = [
                f"**Content of: {fname}** ({n_returned} of {n_parents} sections most relevant to: \"{query.strip()}\")",
                "",
            ]
            if not any_match:
                lines.append("_(No strong keyword matches — showing opening sections)_")
                lines.append("")

            for pidx in parent_indices:
                lines.append(f"[Section {pidx + 1}/{n_parents}]")
                lines.append(parent_chunks[pidx][1])
                lines.append("")

            await self._emit(__event_emitter__, "Done", done=True)
            return "\n".join(lines).strip()

        await self._emit(__event_emitter__, "Done", done=True)
        return result

